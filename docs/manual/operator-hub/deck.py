"""Build the Operator Hub documentation deck (PPTX) from the annotated shots.

Layout rules
------------
* wide screenshots  -> full-width image under the title, legend in columns below
* squarish / tall   -> legend column on the left, image on the right
* very tall strips  -> sliced into side-by-side panels first, so a drawer that
  scrolls for 1500 px still lands on the slide at a readable size
"""
from __future__ import annotations

import pathlib

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
PREPPED = HERE / "_build"
PREPPED.mkdir(exist_ok=True)

OUT_DIR = HERE.parent
OUT_PPTX = OUT_DIR / "operator-hub-guide.pptx"

# ---- palette (matches docs/manual/build_operator_manual.py) ----------------
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
ORANGE_LT = RGBColor(0xFF, 0xF3, 0xEB)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LT = RGBColor(0xF3, 0xF4, 0xF6)
BLUE = RGBColor(0x1E, 0x6F, 0xD9)
GREEN = RGBColor(0x0E, 0x8A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
FONT = "Segoe UI"
MONO = "Consolas"

COLOR_BY_NAME = {"": ORANGE, "blue": BLUE, "green": GREEN}


# ---------------------------------------------------------------- helpers --
def slice_tall(path: pathlib.Path) -> pathlib.Path:
    """Cut a tall narrow strip into side-by-side panels so it fits a slide."""
    img = Image.open(path)
    w, h = img.size
    if h / w < 1.7:
        return path
    n = 2
    while (h / n) / (w * n) > 1.15 and n < 4:
        n += 1
    overlap = int(h * 0.012)
    part_h = h // n + overlap
    gap = 26
    out = Image.new("RGB", (w * n + gap * (n - 1), part_h), (223, 226, 230))
    for i in range(n):
        top = min(h - part_h, i * (h // n))
        out.paste(img.crop((0, top, w, top + part_h)), (i * (w + gap), 0))
    dest = PREPPED / f"{path.stem}-split.png"
    out.save(dest, optimize=True)
    return dest


def txbox(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, font=FONT, space_after=4, line=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, part in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        run = p.add_run()
        run.text = part
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def disc(slide, x, y, d, number, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT
    return shape


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_chrome(slide, eyebrow, title, subtitle=""):
    rect(slide, 0, 0, SW, Inches(0.13), ORANGE)
    y = Inches(0.36)
    if eyebrow:
        txbox(slide, MARGIN, y, Inches(11.5), Inches(0.24),
              eyebrow.upper(), size=11, bold=True, color=ORANGE)
        y += Inches(0.30)
    txbox(slide, MARGIN, y, Inches(12.2), Inches(0.5), title, size=25, bold=True)
    y += Inches(0.52)
    if subtitle:
        txbox(slide, MARGIN, y, Inches(12.2), Inches(0.34), subtitle,
              size=12.5, color=GRAY)
        y += Inches(0.38)
    return y + Inches(0.10)


def est_height(text, w_emu, size, line=1.18):
    """Rendered height of wrapped text. Segoe UI averages ~0.5 em per glyph."""
    w_pt = w_emu / 12700.0
    # 0.50 em average glyph width, then 0.86 for the ragged right edge that
    # word wrapping leaves behind.
    per_line = max(8, int(w_pt / (size * 0.50) * 0.86))
    lines = max(1, -(-len(text) // per_line))
    # PowerPoint's single line spacing is ~1.2x the point size; `line` is a
    # multiplier on top of that.
    return Pt(size * 1.2 * line * lines)


def legend_block(slide, items, x, y, w, size=11.5, gap=Inches(0.085)):
    """Numbered callout legend. Returns the y coordinate just past the block."""
    d = Inches(0.235)
    for item in items:
        n, text = item[0], item[1]
        color = COLOR_BY_NAME[item[2] if len(item) > 2 else ""]
        tx, tw = (x + d + Inches(0.11), w - d - Inches(0.11)) if n else (x, w)
        h = est_height(text, tw, size)
        if n:
            disc(slide, x, y + Inches(0.015), d, n, color)
        txbox(slide, tx, y, tw, h, text, size=size, line=1.18)
        y += h + gap
    return y


def legend_height(items, w, size, gap=Inches(0.085)):
    d = Inches(0.235)
    total = 0
    for item in items:
        tw = w - d - Inches(0.11) if item[0] else w
        total += est_height(item[1], tw, size) + gap
    return total


def add_image(slide, path, x, y, max_w, max_h):
    img = Image.open(path)
    ratio = img.width / img.height
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        h = max_h
        w = int(h * ratio)
    left = x + int((max_w - w) / 2)
    top = y + int((max_h - h) / 2)
    frame = rect(slide, left - Emu(9000), top - Emu(9000),
                 w + Emu(18000), h + Emu(18000), WHITE, line=RGBColor(0xC9, 0xCD, 0xD3))
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)
    return frame


def note_bar(slide, y, text, color=BLUE):
    tint = RGBColor(0xEF, 0xF6, 0xFF) if color == BLUE else ORANGE_LT
    h = Inches(0.46)
    rect(slide, MARGIN, y, SW - 2 * MARGIN, h, tint)
    rect(slide, MARGIN, y, Inches(0.045), h, color)
    txbox(slide, MARGIN + Inches(0.20), y + Inches(0.11),
          SW - 2 * MARGIN - Inches(0.4), Inches(0.3), text, size=11, color=DARK)


# ------------------------------------------------------------ slide kinds --
def slide_cover(prs, title, subtitle, meta):
    slide = blank(prs)
    rect(slide, 0, 0, SW, SH, DARK)
    rect(slide, 0, 0, Inches(0.30), SH, ORANGE)
    txbox(slide, Inches(1.1), Inches(2.05), Inches(11), Inches(0.35),
          "PERTAMINA GLD", size=15, bold=True, color=ORANGE)
    txbox(slide, Inches(1.1), Inches(2.55), Inches(11.2), Inches(1.5),
          title, size=44, bold=True, color=WHITE, line=1.05)
    rect(slide, Inches(1.1), Inches(4.35), Inches(2.0), Inches(0.05), ORANGE)
    txbox(slide, Inches(1.1), Inches(4.70), Inches(10.5), Inches(0.9),
          subtitle, size=16, color=RGBColor(0xC9, 0xCD, 0xD8), line=1.35)
    txbox(slide, Inches(1.1), Inches(6.35), Inches(10.5), Inches(0.5),
          meta, size=11, color=GRAY, line=1.3)
    return slide


def slide_section(prs, number, title, subtitle, bullets):
    slide = blank(prs)
    rect(slide, 0, 0, SW, SH, DARK)
    rect(slide, 0, 0, Inches(0.30), SH, ORANGE)
    txbox(slide, Inches(1.1), Inches(1.55), Inches(9), Inches(0.5),
          f"BAGIAN {number}", size=14, bold=True, color=ORANGE)
    txbox(slide, Inches(1.1), Inches(2.05), Inches(11), Inches(1.1),
          title, size=40, bold=True, color=WHITE)
    txbox(slide, Inches(1.1), Inches(3.30), Inches(10.5), Inches(0.6),
          subtitle, size=15, color=RGBColor(0xC9, 0xCD, 0xD8))
    y = Inches(4.20)
    for i, b in enumerate(bullets, 1):
        disc(slide, Inches(1.1), y + Inches(0.02), Inches(0.26), i, ORANGE)
        txbox(slide, Inches(1.55), y, Inches(10.6), Inches(0.3), b,
              size=13, color=WHITE)
        y += Inches(0.42)
    return slide


def slide_text(prs, eyebrow, title, subtitle, blocks, note=None):
    """blocks: list of (heading, [lines]) rendered as up to three columns."""
    slide = blank(prs)
    top = page_chrome(slide, eyebrow, title, subtitle)
    cols = min(3, max(1, len(blocks)))
    gap = Inches(0.35)
    cw = int((SW - 2 * MARGIN - gap * (cols - 1)) / cols)
    for i, (heading, lines) in enumerate(blocks):
        x = MARGIN + i % cols * (cw + gap)
        y = top + (i // cols) * Inches(2.85)
        rect(slide, x, y, cw, Inches(0.045), ORANGE)
        txbox(slide, x, y + Inches(0.18), cw, Inches(0.3), heading,
              size=14.5, bold=True)
        yy = y + Inches(0.62)
        size = 11.5
        for line in lines:
            if line.startswith("• "):
                tw = cw - Inches(0.18)
                h = est_height(line[2:], tw, size, line=1.2)
                txbox(slide, x + Inches(0.18), yy, tw, h, line[2:],
                      size=size, line=1.2)
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.02),
                                             yy + Inches(0.065),
                                             Inches(0.075), Inches(0.075))
                dot.fill.solid()
                dot.fill.fore_color.rgb = ORANGE
                dot.line.fill.background()
                dot.shadow.inherit = False
            else:
                h = est_height(line, cw, size, line=1.2) if line else Pt(6)
                if line:
                    txbox(slide, x, yy, cw, h, line, size=size, line=1.2)
            yy += h + Inches(0.055)
    if note:
        note_bar(slide, Inches(6.70), note)
    return slide


def slide_shot(prs, eyebrow, title, subtitle, image, legend, note=None,
               force_side=False):
    slide = blank(prs)
    top = page_chrome(slide, eyebrow, title, subtitle)
    path = slice_tall(SHOTS / f"{image}.png")
    img = Image.open(path)
    ratio = img.width / img.height
    note_h = Inches(0.62) if note else Inches(0)

    if ratio >= 1.85 and not force_side:
        # wide: image across the slide, legend underneath in columns
        cols = 3 if len(legend) > 4 else max(1, len(legend))
        cw = int((SW - 2 * MARGIN - Inches(0.3) * (cols - 1)) / cols)
        per = -(-len(legend) // cols)
        size = 11
        legend_h = max(legend_height(legend[c * per:(c + 1) * per], cw, size)
                       for c in range(cols))
        img_h = SH - top - legend_h - note_h - Inches(0.45)
        add_image(slide, path, MARGIN, top, SW - 2 * MARGIN, img_h)
        ly = top + img_h + Inches(0.24)
        for c in range(cols):
            legend_block(slide, legend[c * per:(c + 1) * per],
                         MARGIN + c * (cw + Inches(0.3)), ly, cw, size=size)
    else:
        # tall/square: legend on the left, image on the right
        lw = Inches(3.65)
        size = 11.5
        avail = SH - top - note_h - Inches(0.30)
        while legend_height(legend, lw, size) > avail and size > 9:
            size -= 0.5
        legend_block(slide, legend, MARGIN, top + Inches(0.05), lw, size=size)
        ix = MARGIN + lw + Inches(0.45)
        add_image(slide, path, ix, top, SW - ix - MARGIN, avail)
    if note:
        note_bar(slide, SH - Inches(0.30) - Inches(0.46), note)
    return slide
